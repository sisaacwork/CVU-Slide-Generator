"""
pptx_gen.py — injects query data into the CVU slide template.

Strategy:
  1. python-pptx for text placeholder replacement.
  2. zipfile manipulation to update chart XML caches + embedded Excel files.
  3. Returns PPTX bytes ready to stream as a download.
"""

import io
import zipfile
import datetime
import logging
from pathlib import Path

from pptx import Presentation
import openpyxl
from lxml import etree

log = logging.getLogger(__name__)

TEMPLATE_PATH = Path(__file__).parent / "template.pptx"

NS = {
    'c': 'http://schemas.openxmlformats.org/drawingml/2006/chart',
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006',
}


# ─── Text replacement ────────────────────────────────────────────────────────

def _replace_in_pptx(pptx_bytes: bytes, replacements: dict) -> bytes:
    prs = Presentation(io.BytesIO(pptx_bytes))
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                full = ''.join(run.text for run in para.runs)
                changed = full
                for key, val in replacements.items():
                    changed = changed.replace(key, val)
                if changed != full and para.runs:
                    para.runs[0].text = changed
                    for run in para.runs[1:]:
                        run.text = ''
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ─── Chart XML helpers ───────────────────────────────────────────────────────

def _make_str_cache(cats, ns_c):
    sc = etree.Element(f'{{{ns_c}}}strCache')
    pc = etree.SubElement(sc, f'{{{ns_c}}}ptCount')
    pc.set('val', str(len(cats)))
    for i, cat in enumerate(cats):
        pt = etree.SubElement(sc, f'{{{ns_c}}}pt')
        pt.set('idx', str(i))
        etree.SubElement(pt, f'{{{ns_c}}}v').text = str(cat)
    return sc


def _make_num_cache(vals, fmt, ns_c):
    nc = etree.Element(f'{{{ns_c}}}numCache')
    etree.SubElement(nc, f'{{{ns_c}}}formatCode').text = fmt
    non_none = [(i, v) for i, v in enumerate(vals) if v is not None]
    etree.SubElement(nc, f'{{{ns_c}}}ptCount').set('val', str(len(non_none)))
    for i, v in non_none:
        pt = etree.SubElement(nc, f'{{{ns_c}}}pt')
        pt.set('idx', str(i))
        etree.SubElement(pt, f'{{{ns_c}}}v').text = str(v)
    return nc


def _update_formula_range(formula: str, count: int) -> str:
    import re
    return re.sub(r'(\$[A-Z]+\$)\d+$', lambda m: f"{m.group(1)}{count + 1}", formula)


def _update_chart_series(chart_xml: bytes, series_list: list) -> bytes:
    ns_c = NS['c']
    root = etree.fromstring(chart_xml)
    all_sers = root.findall(f'.//{{{ns_c}}}ser')

    for spec in series_list:
        idx = spec['idx']
        if idx >= len(all_sers):
            continue
        ser   = all_sers[idx]
        count = len(spec.get('cats', spec.get('vals', [])))

        if 'cats' in spec:
            str_ref = ser.find(f'.//{{{ns_c}}}strRef')
            if str_ref is not None:
                f_el = str_ref.find(f'{{{ns_c}}}f')
                if f_el is not None:
                    f_el.text = _update_formula_range(f_el.text or '', len(spec['cats']))
                old = str_ref.find(f'{{{ns_c}}}strCache')
                if old is not None:
                    str_ref.remove(old)
                str_ref.append(_make_str_cache(spec['cats'], ns_c))

        if 'vals' in spec:
            fmt     = spec.get('val_fmt', 'General')
            num_ref = ser.find(f'.//{{{ns_c}}}numRef')
            if num_ref is not None:
                f_el = num_ref.find(f'{{{ns_c}}}f')
                if f_el is not None:
                    f_el.text = _update_formula_range(f_el.text or '', count)
                old = num_ref.find(f'{{{ns_c}}}numCache')
                if old is not None:
                    num_ref.remove(old)
                num_ref.append(_make_num_cache(spec['vals'], fmt, ns_c))

    return etree.tostring(root, xml_declaration=True, encoding='UTF-8')


# ─── Excel helpers ───────────────────────────────────────────────────────────

def _build_excel(headers: list, rows: list) -> bytes:
    """Create a fresh workbook — avoids any corruption from original xlsx."""
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = 'Sheet1'
    ws.append(headers)
    for row in rows:
        ws.append(list(row))
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ─── Per-chart builders ──────────────────────────────────────────────────────

def _chart1(s2):
    years, vals = s2['years'], s2['values']
    series  = [{'idx': 0, 'cats': years, 'vals': vals}]
    xl_rows = list(zip(years, vals))
    return series, ['Year', 'Total Buildings'], xl_rows

def _chart2(s3):
    years, cv, ov = s3['years'], s3['city_values'], s3['other_values']
    series  = [{'idx': 0, 'cats': years, 'vals': cv},
               {'idx': 1, 'cats': years, 'vals': ov}]
    xl_rows = list(zip(years, cv, ov))
    return series, ['Year', 'City', 'All Other Cities'], xl_rows

def _chart3(s4):
    years, hv, pv = s4['years'], s4['hist_values'], s4['proj_values']
    series  = [{'idx': 0, 'cats': years, 'vals': hv},
               {'idx': 1, 'cats': years, 'vals': pv}]
    xl_rows = list(zip(years, hv, pv))
    return series, ['Year', 'Total Buildings', 'Projected'], xl_rows

def _chart4(s5):
    years = s5['years']
    series = [
        {'idx': 0, 'cats': years, 'vals': s5['hist_builds']},
        {'idx': 1, 'cats': years, 'vals': s5['proj_builds']},
        {'idx': 2, 'cats': years, 'vals': s5['hist_pop'],   'val_fmt': '#,##0'},
        {'idx': 3, 'cats': years, 'vals': s5['future_pop'], 'val_fmt': '#,##0'},
    ]
    xl_rows = list(zip(years, s5['hist_builds'], s5['proj_builds'],
                       s5['hist_pop'], s5['future_pop']))
    return series, ['Year', 'Total Buildings', 'Future Buildings',
                    'Total Population', 'Future Population'], xl_rows

def _chart5(s6):
    cats, vals = s6['func_categories'], s6['func_values']
    series  = [{'idx': 0, 'cats': cats, 'vals': vals}]
    xl_rows = list(zip(cats, vals))
    return series, [' ', 'Building Function'], xl_rows

def _chart6(s6):
    cats, vals = s6['mat_categories'], s6['mat_values']
    series  = [{'idx': 0, 'cats': cats, 'vals': vals}]
    xl_rows = list(zip(cats, vals))
    return series, [' ', 'Structural Material'], xl_rows


# ─── Main entry point ────────────────────────────────────────────────────────

def generate_pptx(
    city_name:       str,
    country_name:    str,
    threshold:       int,
    slide_data:      dict,
    selected_slides: list | None = None,
) -> bytes:
    today      = datetime.date.today()
    month_abbr = today.strftime('%b').upper()

    # Cover slide: handle template that has {city}{country} with no separator
    city_country = (f'{city_name}, {country_name}' if country_name
                    else city_name)

    replacements = {
        # Combined patterns first — catches no-separator template layouts
        '{city}{country}':   city_country,
        '{city}, {country}': city_country,
        # Individual fallbacks
        '{city}':            city_name,
        '{country}':         country_name,
        '{threshold}':       str(threshold),
        '{year}':            str(today.year),
        '{month}':           month_abbr,
        '{day}':             f"{today.day:02d}",
        '{2000_growth_pct}': str((slide_data.get('s2') or {}).get('growth_pct_2000', 0)),
    }

    template_bytes = TEMPLATE_PATH.read_bytes()
    pptx_bytes     = _replace_in_pptx(template_bytes, replacements)

    in_zip  = zipfile.ZipFile(io.BytesIO(pptx_bytes))
    out_buf = io.BytesIO()

    # Map chart filename → (series_list, excel_headers, excel_rows, excel_fname)
    chart_map: dict[str, tuple] = {}

    s2 = slide_data.get('s2')
    s3 = slide_data.get('s3')
    s4 = slide_data.get('s4')
    s5 = slide_data.get('s5')
    s6 = slide_data.get('s6')

    if s2:
        ser, hdrs, rows = _chart1(s2)
        chart_map['chart1.xml'] = (ser, hdrs, rows, 'Microsoft_Excel_Worksheet.xlsx')
    if s3:
        ser, hdrs, rows = _chart2(s3)
        chart_map['chart2.xml'] = (ser, hdrs, rows, 'Microsoft_Excel_Worksheet1.xlsx')
    if s4:
        ser, hdrs, rows = _chart3(s4)
        chart_map['chart3.xml'] = (ser, hdrs, rows, 'Microsoft_Excel_Worksheet2.xlsx')
    if s5:
        ser, hdrs, rows = _chart4(s5)
        chart_map['chart4.xml'] = (ser, hdrs, rows, 'Microsoft_Excel_Worksheet3.xlsx')
    if s6:
        ser5, hdrs5, rows5 = _chart5(s6)
        chart_map['chart5.xml'] = (ser5, hdrs5, rows5, 'Microsoft_Excel_Worksheet4.xlsx')
        ser6, hdrs6, rows6 = _chart6(s6)
        chart_map['chart6.xml'] = (ser6, hdrs6, rows6, 'Microsoft_Excel_Worksheet5.xlsx')

    # Pre-build all Excel files (fresh workbooks — avoids corruption issues)
    excel_cache: dict[str, bytes] = {}
    for chart_name, (_, hdrs, rows, xl_fname) in chart_map.items():
        if xl_fname not in excel_cache:
            try:
                excel_cache[xl_fname] = _build_excel(hdrs, rows)
            except Exception as e:
                log.warning("Excel build failed for %s: %s", xl_fname, e)

    with zipfile.ZipFile(out_buf, 'w', zipfile.ZIP_DEFLATED) as out_zip:
        for item in in_zip.infolist():
            content = in_zip.read(item.filename)
            fname   = item.filename

            if 'ppt/charts/' in fname and fname.endswith('.xml'):
                chart_basename = Path(fname).name
                if chart_basename in chart_map:
                    try:
                        content = _update_chart_series(content, chart_map[chart_basename][0])
                    except Exception as e:
                        log.warning("Chart XML update failed for %s: %s", fname, e)

            elif 'ppt/embeddings/' in fname and fname.endswith('.xlsx'):
                xl_basename = Path(fname).name
                if xl_basename in excel_cache:
                    content = excel_cache[xl_basename]

            out_zip.writestr(item, content)

    in_zip.close()
    return out_buf.getvalue()